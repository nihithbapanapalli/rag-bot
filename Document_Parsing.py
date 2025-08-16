import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
import numpy as np
import faiss
from sentence_transformers import SentenceTransformer
from openai import OpenAI
client = OpenAI(
  base_url="https://openrouter.ai/api/v1",
  api_key="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",  # Replace with your API key
)

# Initialize FAISS and Sentence Transformer
embedder = SentenceTransformer("all-MiniLM-L6-v2")
dimension = 384
faiss_index = faiss.IndexFlatL2(dimension)
document_texts = []  # Stores extracted text


def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    extracted_text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():  # Check if shape contains text
                extracted_text.append(shape.text.strip())

    return "\n".join(extracted_text)

# Function to extract text based on file extension
def extract_text_from_document(file_path):
    file_extension = os.path.splitext(file_path)[1].lower()

    if file_extension == ".pdf":
        reader = PdfReader(file_path)
        return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
    elif file_extension == ".docx":
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    elif file_extension == ".pptx":
        return extract_text_from_pptx(file_path)
    elif file_extension == ".html":
        with open(file_path, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f, "html.parser")
            return soup.get_text().strip()
    else:
        return "[ERROR] Unsupported file format."

# Function to add extracted text to FAISS
def add_document_to_faiss(text):
    global faiss_index, document_texts
    sentences = [s.strip() for s in text.split("\n") if s.strip()]
    batch_size = 32

    for i in range(0, len(sentences), batch_size):
        batch = sentences[i: i + batch_size]
        embeddings = embedder.encode(batch)
        embeddings = np.array(embeddings, dtype=np.float32)
        faiss.normalize_L2(embeddings)
        faiss_index.add(embeddings)
        document_texts.extend(batch)

# Function to process a new file
def upload_document(file_path):
    extracted_text = extract_text_from_document(file_path)
    if extracted_text.startswith("[ERROR]"):
        print("Unsupported file format!")
        return
    add_document_to_faiss(extracted_text)
    print("Document added successfully!")

# Search function
def search_document(query, index, doc_texts):
    query_vec = embedder.encode(query).astype('float32')
    _, I = index.search(np.array([query_vec]), k=10)  # Return top 3 results
    return [doc_texts[i] for i in I[0]]

def generate_answer(query, index, doc_texts):
    retrieved_texts = search_document(query, index, doc_texts)
    context = " ".join(retrieved_texts)
    prompt = f"""you are a top researcher trying to read the documents and form a proper structure  only by using the Context: {context} to properly answer the question {query}
              in 50 to 70 words"""
    completion = client.chat.completions.create(
        model="mistralai/mistral-7b-instruct:free",
        messages=[{"role": "user", "content": prompt}]
    )
    return completion.choices[0].message.content


def summarize_document():
    if not document_texts:
        return "[ERROR] No documents have been uploaded yet."

    document_content = " ".join(document_texts)
    print(document_content)# Convert list to a single text block

    prompt = f"Summarize the following document in simple terms:\n{document_content}\nSummary:"

    completion = client.chat.completions.create(
        model="mistralai/mistral-7b-instruct:free",
        messages=[{"role": "user", "content": prompt}]
    )

    return completion.choices[0].message.content



